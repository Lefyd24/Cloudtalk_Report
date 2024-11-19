from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime as dt
import logging
import time
from warnings import filterwarnings
from modified_report_script import ETL
from send_email import send_email

filterwarnings('ignore')
logging.basicConfig(filename='.log/run.log', level=logging.INFO,
                    format='%(asctime)s - %(levelname)s - %(message)s')


def add_dataframe_to_slide(prs, title, dataframe, max_rows=15):
    """Add a table to a PowerPoint slide for a given dataframe, splitting if needed."""
    num_slices = (len(dataframe) // max_rows) + 1  # Determine the number of slides needed

    for i in range(num_slices):
        # Create a new slide for each slice of the dataframe
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use the title + content layout
        slide.shapes.title.text = f"{title} (Page {i + 1})"

        # Subset the dataframe for the current slide
        subset = dataframe.iloc[i * max_rows:(i + 1) * max_rows]

        # Define table dimensions
        rows, cols = subset.shape
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.2 * (rows + 1))
        
        # Adjust table creation to handle column limits dynamically
        table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

        # Style the table headers
        for j, column_name in enumerate(subset.columns):
            cell = table.cell(0, j)
            cell.text = column_name.replace('_', ' ').capitalize()
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Light blue for headers
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add table data
        for row_idx, row in subset.iterrows():
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(round(value, 2)) if isinstance(value, float) else str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Set date range
date_to = dt.datetime.now().strftime('%Y-%m-%d %H:%M:%S')
date_from = (dt.datetime.now() - dt.timedelta(days=7)).strftime('%Y-%m-%d 00:00:00')

logging.info(f'Generating report for the date range: {date_from} - {date_to}')
start = time.time()

try:
    answered_df, missed_df, error_calls, agents, agents_ext, unanswered_call_counts, unanswered_calls_ring, unanswered_calls_hour = ETL(date_from, date_to)
except Exception as e:
    print(e)
    logging.error(f'Error while generating report: {e}')
    exit(1)

# Create PowerPoint presentation
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'CloudTalk - Mitsis Group Analytics'
subtitle = slide.placeholders[1]
subtitle.text = f"Date range: {dt.datetime.strptime(date_from, '%Y-%m-%d %H:%M:%S').strftime('%d/%m/%Y %H:%M:%S')} - {dt.datetime.strptime(date_to, '%Y-%m-%d %H:%M:%S').strftime('%d/%m/%Y %H:%M:%S')}"

# General statistics slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'General Statistics'
content = slide.shapes.placeholders[1].text_frame
content.text = f"Total Answered calls: {len(answered_df)}"
content.add_paragraph().text = f"Total Missed calls: {len(missed_df)}"
content.add_paragraph().text = f"Total Calls that failed to be retrieved through API: {error_calls}"

# Add tables
add_dataframe_to_slide(prs, '1. Agent Statistics (incl. Internal Calls)', agents)
add_dataframe_to_slide(prs, '2. Agent Statistics (excl. Internal Calls)', agents_ext)
add_dataframe_to_slide(prs, '3. Unanswered Calls per Agent', unanswered_call_counts)

# Save the presentation
prs.save('Cloudtalk_Report_Pretty.pptx')

logging.info(f'Report generated successfully in {round(time.time() - start, 2)} seconds')